from pptx import Presentation
from transformers import pipeline
import os
import re

# Extract text and images from .pptx file
def extract_text_and_images_from_pptx(pptx_file, output_image_dir):
    prs = Presentation(pptx_file)
    slides_content = []
    
    # Ensure image output directory exists
    os.makedirs(output_image_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides):
        slide_content = {"text": "", "images": []}
        
        # Extract slide title
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title = shape.text.strip().split("\n")[0]  # Take the first line as the title
                break

        slide_content["title"] = title if title else f"Slide {i+1}"
        
        # Extract text and clean it
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Clean and replace weird encoding characters (e.g., diamonds with ? inside)
                text = shape.text.strip().replace("�", "")  # Remove unexpected characters
                text = re.sub(r'[^\x00-\x7F]+', '', text)  # Remove non-ASCII characters
                
                # Break long paragraphs into bullet points by splitting sentences
                bullet_points = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)
                slide_content["text"] += "\n".join([f"- {point.strip()}" for point in bullet_points if point.strip()]) + "\n"

        # Extract images
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_filename = os.path.join(output_image_dir, f"slide_{i+1}_image_{len(slide_content['images'])+1}.png")
                with open(image_filename, "wb") as img_file:
                    img_file.write(image.blob)
                slide_content["images"].append(image_filename)

        slides_content.append(slide_content)

    return slides_content

# Summarize each slide using Hugging Face model
def summarize_slides(slides_content):
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    summarized_notes = []
    
    for slide in slides_content:
        text = slide["text"].strip()
        if len(text) > 0:  # Summarize only non-empty text
            summary = summarizer(text, max_length=250, min_length=80, do_sample=False)[0]['summary_text']
            # Format the summary as bullet points
            summary_bullets = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', summary)
            bullet_summary = "\n".join([f"- {s.strip()}" for s in summary_bullets if s.strip()])
        else:
            bullet_summary = "(No significant text on this slide)"
        summarized_notes.append({"title": slide["title"], "summary": bullet_summary, "images": slide["images"]})
    
    return summarized_notes

# Save summarized notes and image links to markdown for Obsidian
def save_as_markdown(summarized_notes, output_file, image_dir):
    with open(output_file, "w") as f:
        for i, slide in enumerate(summarized_notes):
            f.write(f"# {slide['title']}\n")  # Use the title extracted from the slide
            f.write(f"{slide['summary']}\n\n")
            
            # Link images in markdown using HTML for resizing
            for img_path in slide["images"]:
                relative_path = os.path.relpath(img_path, image_dir)
                # Resize using HTML <img> tag for Obsidian compatibility
                f.write(f'<img src="{relative_path}" width="500px">\n')
            f.write("\n")

# Main function to tie everything together
def process_pptx_to_detailed_notes(pptx_file, output_file, output_image_dir):
    slides_content = extract_text_and_images_from_pptx(pptx_file, output_image_dir)
    summarized_notes = summarize_slides(slides_content)
    save_as_markdown(summarized_notes, output_file, output_image_dir)
    print(f"Summarized notes with images saved to {output_file}")

# Update with your .pptx file path, output markdown file, and image directory
pptx_file = "W1L1-comp-org-intro.pptx"  # Replace with your file path
output_file = "notes_for_obsidian.md"  # Replace with your desired output file name
output_image_dir = r"C:\Users\tom\Documents\LLM_Playground\photos"  # Directory to save extracted images

process_pptx_to_detailed_notes(pptx_file, output_file, output_image_dir)
